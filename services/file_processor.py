"""
Enhanced File Processing Service for handling uploaded files with content extraction
"""

import os
import re
import mimetypes
import base64
from typing import Dict, Any, Optional, Tuple, Union
from fastapi import UploadFile
import json
from datetime import datetime

# Import document processing libraries
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class FileProcessor:
    """Service for processing uploaded files"""
    
    SUPPORTED_TEXT_TYPES = {
        'text/plain',
        'text/markdown',
        'text/csv',
        'application/json',
        'application/xml',
        'text/xml',
        'application/yaml',
        'text/yaml',
        'application/x-yaml',
        'text/html',
        'text/css',
        'text/javascript',
        'application/javascript',
    }
    
    SUPPORTED_DOCUMENT_TYPES = {
        'application/pdf',
        'application/msword',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'application/vnd.ms-excel',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'application/vnd.ms-powerpoint',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    }
    
    SUPPORTED_FILE_TYPES = SUPPORTED_TEXT_TYPES.union(SUPPORTED_DOCUMENT_TYPES)
    
    @classmethod
    def get_file_info(cls, file: UploadFile) -> Dict[str, Any]:
        """Get file information"""
        return {
            "filename": file.filename,
            "content_type": file.content_type,
            "size": getattr(file, 'size', None),
        }
    
    @classmethod
    def is_supported_file_type(cls, content_type: str) -> bool:
        """Check if file type is supported"""
        return content_type in cls.SUPPORTED_FILE_TYPES
    
    @classmethod
    def is_text_file(cls, content_type: str) -> bool:
        """Check if file is a text-based file"""
        return content_type in cls.SUPPORTED_TEXT_TYPES
    
    @classmethod
    async def read_file_content(cls, file: UploadFile) -> Tuple[str, str]:
        """Read file content with enhanced extraction for different file types"""
        content = await file.read()
        filename = file.filename or "unknown"
        content_type = file.content_type or "application/octet-stream"
        
        # Handle text files
        if cls.is_text_file(content_type):
            try:
                text_content = content.decode('utf-8')
                return text_content, filename
            except UnicodeDecodeError:
                # Try other encodings
                for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                    try:
                        text_content = content.decode(encoding)
                        return text_content, filename
                    except UnicodeDecodeError:
                        continue
                # If all fail, return as base64
                encoded_content = base64.b64encode(content).decode('utf-8')
                return f"Binary file content (base64): {encoded_content}", filename
        
        # Handle document files
        elif content_type in cls.SUPPORTED_DOCUMENT_TYPES:
            extracted_content = cls.extract_document_content(content, content_type, filename)
            return extracted_content, filename
        
        # Fallback for unsupported files
        else:
            encoded_content = base64.b64encode(content).decode('utf-8')
            return f"Unsupported file type: {content_type}\nBinary content (base64): {encoded_content}", filename

    @classmethod
    def extract_document_content(cls, content: bytes, content_type: str, filename: str) -> str:
        """Extract text content from various document formats"""
        try:
            if content_type == 'application/pdf':
                return cls.extract_pdf_content(content)
            elif content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                return cls.extract_docx_content(content)
            elif content_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                return cls.extract_xlsx_content(content)
            elif content_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                return cls.extract_pptx_content(content)
            elif content_type == 'application/msword':
                return cls.extract_doc_content(content)
            else:
                return f"Document type {content_type} not yet supported for content extraction"
        except Exception as e:
            return f"Error extracting content from {filename}: {str(e)}"

    @classmethod
    def extract_pdf_content(cls, content: bytes) -> str:
        """Extract text content from PDF"""
        if not PDF_AVAILABLE:
            return "PDF processing not available. Please install PyPDF2."
        
        try:
            from io import BytesIO
            pdf_file = BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_content = []
            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                if page_text.strip():
                    text_content.append(f"--- Page {page_num} ---\n{page_text}\n")
            
            return "\n".join(text_content) if text_content else "No text content found in PDF"
        except Exception as e:
            return f"Error extracting PDF content: {str(e)}"

    @classmethod
    def extract_docx_content(cls, content: bytes) -> str:
        """Extract text content from DOCX"""
        if not DOCX_AVAILABLE:
            return "DOCX processing not available. Please install python-docx."
        
        try:
            from io import BytesIO
            docx_file = BytesIO(content)
            doc = DocxDocument(docx_file)
            
            paragraphs = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text)
            
            return "\n".join(paragraphs) if paragraphs else "No text content found in DOCX"
        except Exception as e:
            return f"Error extracting DOCX content: {str(e)}"

    @classmethod
    def extract_xlsx_content(cls, content: bytes) -> str:
        """Extract text content from XLSX"""
        if not XLSX_AVAILABLE:
            return "XLSX processing not available. Please install openpyxl."
        
        try:
            from io import BytesIO
            xlsx_file = BytesIO(content)
            workbook = openpyxl.load_workbook(xlsx_file)
            
            content_lines = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content_lines.append(f"--- Sheet: {sheet_name} ---")
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_data):
                        content_lines.append("\t".join(row_data))
                content_lines.append("")
            
            return "\n".join(content_lines) if content_lines else "No content found in XLSX"
        except Exception as e:
            return f"Error extracting XLSX content: {str(e)}"

    @classmethod
    def extract_pptx_content(cls, content: bytes) -> str:
        """Extract text content from PPTX"""
        if not PPTX_AVAILABLE:
            return "PPTX processing not available. Please install python-pptx."
        
        try:
            from io import BytesIO
            pptx_file = BytesIO(content)
            presentation = Presentation(pptx_file)
            
            content_lines = []
            for slide_num, slide in enumerate(presentation.slides, 1):
                content_lines.append(f"--- Slide {slide_num} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_lines.append(shape.text)
                content_lines.append("")
            
            return "\n".join(content_lines) if content_lines else "No text content found in PPTX"
        except Exception as e:
            return f"Error extracting PPTX content: {str(e)}"

    @classmethod
    def extract_doc_content(cls, content: bytes) -> str:
        """Extract text content from legacy DOC files"""
        # Legacy DOC files are more complex to parse
        # For now, return a message indicating limitation
        return "Legacy DOC files (.doc) are not fully supported. Please convert to DOCX format for better content extraction."
    
    @classmethod
    def generate_filename(cls, original_filename: str, content_type: str) -> str:
        """Generate a safe filename"""
        if not original_filename:
            # Generate filename based on content type
            extension = mimetypes.guess_extension(content_type) or '.txt'
            return f"uploaded_file{extension}"
        
        # Clean filename
        safe_filename = re.sub(r'[^\w\s.-]', '', original_filename)
        safe_filename = re.sub(r'[-\s]+', '-', safe_filename)
        
        return safe_filename
    
    @classmethod
    def create_file_metadata(cls, file: UploadFile, content: str) -> Dict[str, Any]:
        """Create metadata for uploaded file (standardized format)"""
        file_info = cls.get_file_info(file)
        
        metadata = {
            "original_filename": file_info["filename"],
            "content_type": file_info["content_type"],
            "file_size": len(content.encode('utf-8')),
            "content": content,
            "extraction_method": cls.get_extraction_method(file_info["content_type"]),
            "upload_timestamp": datetime.utcnow().isoformat(),
            "is_text_file": cls.is_text_file(file_info["content_type"]),
            "is_document_file": file_info["content_type"] in cls.SUPPORTED_DOCUMENT_TYPES,
            "source_type": "file"
        }
        
        return metadata

    @classmethod
    def get_extraction_method(cls, content_type: str) -> str:
        """Get the extraction method used for the file type"""
        if content_type in cls.SUPPORTED_TEXT_TYPES:
            return "text_decode"
        elif content_type == 'application/pdf':
            return "pdf_extract"
        elif content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return "docx_extract"
        elif content_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            return "xlsx_extract"
        elif content_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return "pptx_extract"
        elif content_type == 'application/msword':
            return "doc_extract"
        else:
            return "base64_encode"
    
    @classmethod
    def format_file_content(cls, content: str, filename: str, content_type: str) -> str:
        """Format file content for storage - just return the raw content"""
        return content
